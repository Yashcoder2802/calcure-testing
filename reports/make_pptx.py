"""Generate reports/presentation.pptx — polished v2."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
import os

# ── Palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x11, 0x17)
SURFACE = RGBColor(0x16, 0x1B, 0x22)
CARD2   = RGBColor(0x1C, 0x23, 0x2C)
BORDER  = RGBColor(0x30, 0x36, 0x3D)
ACCENT  = RGBColor(0x58, 0xA6, 0xFF)
GREEN   = RGBColor(0x3F, 0xB9, 0x50)
ORANGE  = RGBColor(0xE3, 0xB3, 0x41)
RED     = RGBColor(0xF8, 0x51, 0x49)
PURPLE  = RGBColor(0xBC, 0x8C, 0xFF)
TEAL    = RGBColor(0x39, 0xD3, 0xC3)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x89, 0x93, 0x9E)
DIM     = RGBColor(0x48, 0x4F, 0x58)

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]

# ── Low-level helpers ─────────────────────────────────────────────────────────

def add_slide():
    sl = prs.slides.add_slide(blank)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return sl

def box(sl, x, y, w, h, color, line_color=None, line_pt=0):
    shp = sl.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_pt)
    else:
        shp.line.fill.background()
    return shp

def card(sl, x, y, w, h, color=None, border=BORDER):
    c = color or SURFACE
    shp = sl.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = c
    shp.line.color.rgb = border
    shp.line.width = Pt(0.75)
    return shp

def txt(sl, text, x, y, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = sl.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def pill(sl, label, x, y, w=Inches(1.5), h=Inches(0.36), color=ACCENT,
         text_color=WHITE, size=11):
    shp = sl.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color

def metric_card(sl, x, y, w, h, value, label, color=ACCENT, sub=None):
    card(sl, x, y, w, h)
    box(sl, x, y, w, Inches(0.05), color)   # top accent bar
    txt(sl, value, x, y + Inches(0.18), w, Inches(0.72),
        size=40, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(sl, label, x, y + Inches(0.88), w, Inches(0.38),
        size=12, color=MUTED, align=PP_ALIGN.CENTER)
    if sub:
        txt(sl, sub, x, y + Inches(1.22), w, Inches(0.3),
            size=10, color=DIM, align=PP_ALIGN.CENTER)

def progress_bar(sl, x, y, w, pct, color=GREEN, bg=BORDER, h=Inches(0.28)):
    box(sl, x, y, w, h, bg)
    fill_w = max(int(w * pct / 100), 1)
    box(sl, x, y, fill_w, h, color)

def bullet_list(sl, items, x, y, w, item_h=Inches(0.55), size=14,
                dot=ACCENT, text_color=WHITE):
    for i, item in enumerate(items):
        yy = y + i * item_h
        box(sl, x, yy + item_h * 0.42, Inches(0.10), Inches(0.10), dot)
        txt(sl, item, x + Inches(0.22), yy, w - Inches(0.22), item_h,
            size=size, color=text_color)

def header(sl, title, sub=None, color=ACCENT):
    box(sl, 0, 0, SW, Inches(0.07), color)
    txt(sl, title, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.72),
        size=30, bold=True, color=WHITE)
    if sub:
        txt(sl, sub, Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.38),
            size=13, color=MUTED)
    box(sl, 0, SH - Inches(0.05), SW, Inches(0.05), color)

def code_box(sl, code_text, x, y, w, h, border_color=BORDER):
    shp = sl.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x14)
    shp.line.color.rgb = border_color
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0xC9, 0xD1, 0xD9)
    return shp

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER  (completely redesigned)
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()

# Left dark panel
box(sl, 0, 0, Inches(8.5), SH, RGBColor(0x0A, 0x0E, 0x14))
# Right accent panel
box(sl, Inches(8.5), 0, Inches(4.833), SH, RGBColor(0x12, 0x1E, 0x35))
# Vertical accent stripe
box(sl, Inches(8.4), 0, Inches(0.10), SH, ACCENT)

# Top-left tiny label
txt(sl, "SOFTWARE TESTING & DESIGN  —  APRIL 2026",
    Inches(0.55), Inches(0.32), Inches(7.5), Inches(0.38),
    size=10, color=MUTED, italic=True)

# Main title block — left side
txt(sl, "Testing", Inches(0.55), Inches(1.2), Inches(7.7), Inches(1.3),
    size=62, bold=True, color=WHITE)
txt(sl, "Calcure", Inches(0.55), Inches(2.35), Inches(7.7), Inches(1.4),
    size=72, bold=True, color=ACCENT)

# Subtitle
txt(sl, "A comprehensive test suite for a\nPython TUI Calendar & Task Manager",
    Inches(0.55), Inches(3.72), Inches(7.5), Inches(0.9),
    size=16, color=MUTED)

# Divider line
box(sl, Inches(0.55), Inches(4.75), Inches(3.5), Inches(0.04), DIM)

# Authors
txt(sl, "Yash Salunke  ·  Harshal Gajjar",
    Inches(0.55), Inches(4.9), Inches(7.5), Inches(0.45),
    size=14, bold=True, color=WHITE)
txt(sl, "github.com/Yashcoder2802/calcure-testing",
    Inches(0.55), Inches(5.38), Inches(7.5), Inches(0.38),
    size=11, color=ACCENT)

# Right panel — 4 big stat cards stacked
stats = [
    ("341", "Total Tests",      GREEN),
    ("89%", "Branch Coverage",  ACCENT),
    ("81.6%","Mutation Score",  ORANGE),
    ("2",   "Bugs Found",       RED),
]
for i, (val, lbl, col) in enumerate(stats):
    y = Inches(0.5) + i * Inches(1.62)
    box(sl, Inches(8.65), y, Inches(4.5), Inches(1.45),
        RGBColor(0x16, 0x24, 0x3A))
    box(sl, Inches(8.65), y, Inches(0.08), Inches(1.45), col)
    txt(sl, val, Inches(8.85), y + Inches(0.12),
        Inches(3.5), Inches(0.82), size=44, bold=True, color=col)
    txt(sl, lbl, Inches(8.85), y + Inches(0.92),
        Inches(3.5), Inches(0.38), size=13, color=MUTED)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Agenda", "What we'll cover today")

items = [
    ("01", "What is Calcure?",            "App overview, architecture & scope",     ACCENT),
    ("02", "Testing Strategy",            "6 suites, techniques & toolchain",       GREEN),
    ("03", "Blackbox Testing",            "EP · BA · EG · CT — 117 tests",          ORANGE),
    ("04", "Whitebox & Coverage",         "86 tests · 89% branch coverage",         PURPLE),
    ("05", "Mutation Testing",            "mutmut · 81.6% kill rate · 784 mutants", RED),
    ("06", "CLI · Mock · Integration",    "65 tests across 3 extra suites",         TEAL),
    ("07", "Bugs Revealed",               "2 real faults with root-cause analysis",  RED),
    ("08", "Results & Lessons Learned",   "Before/after & key takeaways",           GREEN),
]
for i, (num, title, desc, color) in enumerate(items):
    col = i % 2
    row = i // 2
    x = Inches(0.35) + col * Inches(6.55)
    y = Inches(1.42) + row * Inches(1.42)
    card(sl, x, y, Inches(6.2), Inches(1.28))
    box(sl, x, y, Inches(0.08), Inches(1.28), color)
    txt(sl, num, x + Inches(0.22), y + Inches(0.1), Inches(0.65), Inches(0.55),
        size=26, bold=True, color=color)
    txt(sl, title, x + Inches(0.9), y + Inches(0.1), Inches(5.0), Inches(0.45),
        size=15, bold=True, color=WHITE)
    txt(sl, desc, x + Inches(0.9), y + Inches(0.56), Inches(5.0), Inches(0.45),
        size=11, color=MUTED)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS CALCURE
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "What is Calcure?", "Open-source Python TUI calendar — ~5,784 LOC", PURPLE)

# Architecture table left
card(sl, Inches(0.35), Inches(1.35), Inches(5.9), Inches(5.78))
txt(sl, "Architecture", Inches(0.55), Inches(1.48), Inches(5.5), Inches(0.38),
    size=13, bold=True, color=PURPLE)
modules = [
    ("data.py",       "Tasks, events, collections"),
    ("calendars.py",  "Calendar rendering & views"),
    ("loaders.py",    "CSV / ICS file parsing"),
    ("savers.py",     "CSV / ICS file writing"),
    ("screen.py",     "Curses UI, key handling"),
    ("app.py",        "App startup & routing"),
    ("config.py",     "User configuration"),
]
for i, (mod, desc) in enumerate(modules):
    y = Inches(1.95) + i * Inches(0.68)
    bg = CARD2 if i % 2 == 0 else SURFACE
    box(sl, Inches(0.4), y, Inches(5.8), Inches(0.64), bg)
    txt(sl, mod, Inches(0.55), y + Inches(0.12), Inches(1.8), Inches(0.4),
        size=12, bold=True, color=PURPLE)
    txt(sl, desc, Inches(2.45), y + Inches(0.12), Inches(3.6), Inches(0.4),
        size=11, color=MUTED)

# Right side — features + stats
card(sl, Inches(6.55), Inches(1.35), Inches(6.45), Inches(2.6))
txt(sl, "Key Features", Inches(6.75), Inches(1.48), Inches(6.1), Inches(0.38),
    size=13, bold=True, color=ACCENT)
feats = [
    "Terminal UI with Python curses",
    "Events: create · edit · repeat · ICS import",
    "Tasks: priority · subtasks · journals",
    "Persian / Julian calendar support",
    "Timezone-aware multi-view navigation",
]
bullet_list(sl, feats, Inches(6.75), Inches(1.9), Inches(6.1),
            item_h=Inches(0.46), size=12, dot=ACCENT)

# Stats row
for i, (val, lbl, col) in enumerate([
    ("5,784", "Lines of Code", GREEN),
    ("7",     "Modules",       ACCENT),
    ("3.9+",  "Python",        ORANGE),
]):
    x = Inches(6.55) + i * Inches(2.15)
    card(sl, x, Inches(4.1), Inches(2.0), Inches(1.4))
    txt(sl, val, x, Inches(4.22), Inches(2.0), Inches(0.62),
        size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, lbl, x, Inches(4.82), Inches(2.0), Inches(0.35),
        size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Testing scope note
card(sl, Inches(6.55), Inches(5.65), Inches(6.45), Inches(1.48))
txt(sl, "Testing Scope", Inches(6.75), Inches(5.75), Inches(6.1), Inches(0.38),
    size=12, bold=True, color=ACCENT)
txt(sl, "Primary targets: data.py · calendars.py · loaders.py · savers.py\n"
        "Excluded: screen.py (pure curses UI — tested via CLI tier-2 instead)",
    Inches(6.75), Inches(6.15), Inches(6.1), Inches(0.85),
    size=11, color=MUTED)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — STRATEGY OVERVIEW (with tool table)
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Testing Strategy Overview", "6 suites · 341 tests · full toolchain", GREEN)

# Column headers
col_xs = [Inches(0.35), Inches(2.2), Inches(4.85), Inches(6.3), Inches(8.5), Inches(10.6)]
col_ws = [Inches(1.8),  Inches(2.6), Inches(1.4),  Inches(2.1),  Inches(2.0),  Inches(2.5)]
hdrs   = ["Suite", "Techniques", "Tests", "Coverage Goal", "Tool", "Status"]

box(sl, Inches(0.35), Inches(1.32), SW - Inches(0.7), Inches(0.48),
    RGBColor(0x1A, 0x22, 0x2E))
for j, h in enumerate(hdrs):
    txt(sl, h, col_xs[j], Inches(1.36), col_ws[j], Inches(0.38),
        size=11, bold=True, color=MUTED)

rows = [
    ("Blackbox",          "EP · BA · EG · CT",      "117", "Spec compliance",  "pytest",       GREEN,  "✓ 117 pass"),
    ("Whitebox",          "Branch / path coverage",  "86",  "89% branches",    "coverage.py",  ACCENT, "✓ 86 pass"),
    ("Integration",       "Component chains",        "21",  "Data pipeline",   "pytest",       PURPLE, "✓ 21 pass"),
    ("Mock",              "Side-effect isolation",   "17",  "Pure unit",       "unittest.mock",ORANGE, "✓ 17 pass"),
    ("CLI",               "Tier1 + Tier2 PTY",       "27",  "Flag correctness","subprocess+pty",RED,   "✓ 27 pass"),
    ("Mutation-Targeted", "Surviving-mutant focus",  "73",  "81.6% kill rate", "mutmut",       TEAL,   "✓ 73 pass"),
]
for i, (suite, tech, count, goal, tool, color, status) in enumerate(rows):
    y = Inches(1.80) + i * Inches(0.84)
    bg = SURFACE if i % 2 == 0 else BG
    box(sl, Inches(0.35), y, SW - Inches(0.7), Inches(0.80), bg)
    box(sl, Inches(0.35), y, Inches(0.07), Inches(0.80), color)
    txt(sl, suite,  col_xs[0] + Inches(0.12), y+Inches(0.17), col_ws[0], Inches(0.45), size=13, bold=True, color=color)
    txt(sl, tech,   col_xs[1], y+Inches(0.17), col_ws[1], Inches(0.45), size=11, color=WHITE)
    txt(sl, count,  col_xs[2], y+Inches(0.1),  col_ws[2], Inches(0.58), size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(sl, goal,   col_xs[3], y+Inches(0.17), col_ws[3], Inches(0.45), size=11, color=MUTED)
    txt(sl, tool,   col_xs[4], y+Inches(0.17), col_ws[4], Inches(0.45), size=11, color=MUTED)
    pill(sl, status, col_xs[5], y+Inches(0.2), Inches(2.2), Inches(0.38), GREEN, size=10)

# Total bar
box(sl, Inches(0.35), Inches(6.87), SW - Inches(0.7), Inches(0.42),
    RGBColor(0x1A, 0x2A, 0x1A))
txt(sl, "TOTAL", col_xs[0] + Inches(0.12), Inches(6.9), Inches(1.8), Inches(0.38),
    size=12, bold=True, color=WHITE)
txt(sl, "341", col_xs[2], Inches(6.87), col_ws[2], Inches(0.42),
    size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BLACKBOX TECHNIQUES (4 cards detailed)
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Blackbox Testing — 4 Techniques", "117 tests derived from specification only", ORANGE)

techs = [
    ("EP", "Equivalence Partitioning", ACCENT,
     [("Valid name (1–999 chars)", "Accepted"),
      ("Invalid name (0 or 1000+)", "Rejected"),
      ("Priority 1–3",             "Accepted"),
      ("Priority 0 or 4+",         "Rejected")],
     "Divides input space into classes where all values should behave identically."),
    ("BA", "Boundary Analysis", GREEN,
     [("Name length 999",  "Valid ✓"),
      ("Name length 1000", "Invalid ✗"),
      ("Priority = 1",     "Min valid ✓"),
      ("Priority = 3",     "Max valid ✓")],
     "Tests values at the exact boundary between valid and invalid partitions."),
    ("EG", "Error Guessing", RED,
     [('Name = "   " (spaces)', "Should reject ✗ BUG!"),
      ("CSV year = 'abc'",       "Should skip ✗ BUG!"),
      ("Emoji in name",           "Handled gracefully"),
      ("Semicolons in name",      "Escaped correctly")],
     "Guesses common mistakes developers make — found 2 real bugs."),
    ("CT", "Combinatorial Testing", PURPLE,
     [("Priority 1 × done status", "Valid combo"),
      ("Priority 3 × pending",     "Valid combo"),
      ("Repeat × interval pairs",  "6 combos tested"),
      ("Config flag combos",       "4 combos tested")],
     "Tests interactions between parameters that could fail together."),
]
for i, (tag, title, color, table_rows, desc) in enumerate(techs):
    col = i % 2
    row = i // 2
    x = Inches(0.35) + col * Inches(6.55)
    y = Inches(1.35) + row * Inches(2.98)
    card(sl, x, y, Inches(6.2), Inches(2.82))
    box(sl, x, y, SW * 0, Inches(0), color)  # placeholder

    # Tag + title
    pill(sl, tag, x + Inches(0.18), y + Inches(0.14), Inches(0.7), Inches(0.34), color, size=12)
    txt(sl, title, x + Inches(1.05), y + Inches(0.14), Inches(4.9), Inches(0.38),
        size=14, bold=True, color=WHITE)
    txt(sl, desc, x + Inches(0.18), y + Inches(0.55), Inches(5.8), Inches(0.45),
        size=10, color=MUTED, italic=True)

    # mini table
    box(sl, x + Inches(0.18), y + Inches(1.05), Inches(5.8), Inches(0.3),
        RGBColor(0x1A, 0x22, 0x2E))
    txt(sl, "Input", x + Inches(0.25), y + Inches(1.08), Inches(3.0), Inches(0.25),
        size=9, bold=True, color=MUTED)
    txt(sl, "Expected", x + Inches(3.4), y + Inches(1.08), Inches(2.5), Inches(0.25),
        size=9, bold=True, color=MUTED)
    for j, (inp, exp) in enumerate(table_rows):
        yy = y + Inches(1.35) + j * Inches(0.34)
        bg = CARD2 if j % 2 == 0 else SURFACE
        box(sl, x + Inches(0.18), yy, Inches(5.8), Inches(0.32), bg)
        txt(sl, inp, x + Inches(0.25), yy + Inches(0.05), Inches(3.0), Inches(0.25),
            size=9, color=WHITE)
        exp_color = RED if "BUG" in exp else (GREEN if "✓" in exp else MUTED)
        txt(sl, exp, x + Inches(3.4), yy + Inches(0.05), Inches(2.5), Inches(0.25),
            size=9, bold=("BUG" in exp), color=exp_color)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — WHITEBOX COVERAGE MAP
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Whitebox Testing — Branch Coverage", "86 tests targeting specific source branches", ACCENT)

# LEFT — module coverage bars
card(sl, Inches(0.35), Inches(1.35), Inches(6.5), Inches(5.78))
txt(sl, "Branch Coverage by Module", Inches(0.55), Inches(1.48), Inches(6.1), Inches(0.38),
    size=13, bold=True, color=ACCENT)

modules = [
    ("calendars.py", 97, GREEN,  "97 / 100 branches"),
    ("data.py",      94, GREEN,  "94 / 100 branches"),
    ("savers.py",    92, ACCENT, "92 / 100 branches"),
    ("loaders.py",   85, ORANGE, "85 / 100 branches"),
    ("__init__.py",  78, ORANGE, "78 / 100 branches"),
]
for i, (mod, pct, color, detail) in enumerate(modules):
    y = Inches(2.0) + i * Inches(0.98)
    txt(sl, mod, Inches(0.55), y, Inches(2.0), Inches(0.38),
        size=13, bold=True, color=WHITE)
    txt(sl, detail, Inches(2.65), y, Inches(4.1), Inches(0.38),
        size=10, color=MUTED)
    progress_bar(sl, Inches(0.55), y + Inches(0.42), Inches(6.1), pct,
                 color=color, h=Inches(0.3))
    txt(sl, f"{pct}%", Inches(0.65), y + Inches(0.42), Inches(1.0), Inches(0.3),
        size=11, bold=True, color=BG)

# Overall
box(sl, Inches(0.35), Inches(6.6), Inches(6.5), Inches(0.04), BORDER)
txt(sl, "Core Average: 89%", Inches(0.55), Inches(6.68), Inches(4.0), Inches(0.38),
    size=14, bold=True, color=GREEN)

# RIGHT — whitebox test classes table
card(sl, Inches(7.1), Inches(1.35), Inches(5.9), Inches(5.78))
txt(sl, "Key Test Classes", Inches(7.3), Inches(1.48), Inches(5.5), Inches(0.38),
    size=13, bold=True, color=ACCENT)
classes = [
    ("TimerParityTest",    "data.py L234",    "Even/odd second alt."),
    ("CollectionGuard",    "data.py L89",     "999/1000 char guard"),
    ("SubtaskPrefixTest",  "data.py L312",    "'--' vs '----'"),
    ("LeapYearTest",       "calendars.py L78","÷4/÷100/÷400"),
    ("CSVDetectionTest",   "loaders.py L145", "Header vs data row"),
    ("ICSParserTest",      "loaders.py L201", "RRULE + DTSTART"),
    ("PersianCalTest",     "calendars.py L200","Jalali conversion"),
    ("FilterYearTest",     "calendars.py L310","Year boundary"),
    ("MoveTaskTest",       "data.py L401",    "pop + insert order"),
]
box(sl, Inches(7.1), Inches(1.9), Inches(5.9), Inches(0.3),
    RGBColor(0x1A, 0x22, 0x2E))
txt(sl, "Class", Inches(7.25), Inches(1.92), Inches(2.0), Inches(0.25), size=9, bold=True, color=MUTED)
txt(sl, "Location", Inches(9.3), Inches(1.92), Inches(1.7), Inches(0.25), size=9, bold=True, color=MUTED)
txt(sl, "What it tests", Inches(11.0), Inches(1.92), Inches(1.8), Inches(0.25), size=9, bold=True, color=MUTED)
for i, (cls, loc, what) in enumerate(classes):
    y = Inches(2.22) + i * Inches(0.54)
    bg = CARD2 if i % 2 == 0 else SURFACE
    box(sl, Inches(7.15), y, Inches(5.8), Inches(0.50), bg)
    txt(sl, cls,  Inches(7.25), y + Inches(0.07), Inches(2.0), Inches(0.35), size=10, color=PURPLE)
    txt(sl, loc,  Inches(9.3),  y + Inches(0.07), Inches(1.7), Inches(0.35), size=9,  color=MUTED)
    txt(sl, what, Inches(11.0), y + Inches(0.07), Inches(1.8), Inches(0.35), size=9,  color=WHITE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MUTATION TESTING DEEP DIVE
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Mutation Testing Deep Dive", "mutmut v2.4.4 — 784 mutants across data.py + calendars.py", ORANGE)

# Top 4 big metrics
for i, (val, lbl, sub, col) in enumerate([
    ("784",   "Total Mutants",  "generated by mutmut",       ACCENT),
    ("640",   "Killed",         "tests detected the change",  GREEN),
    ("144",   "Survived",       "tests did NOT catch these",  RED),
    ("81.6%", "Kill Rate",      "industry target: ≥80%",      ORANGE),
]):
    metric_card(sl, Inches(0.35 + i * 3.25), Inches(1.35), Inches(3.0), Inches(1.65),
                val, lbl, col, sub)

# Left — surviving mutant analysis
card(sl, Inches(0.35), Inches(3.2), Inches(6.2), Inches(3.93))
txt(sl, "Why Did 144 Mutants Survive?", Inches(0.55), Inches(3.32), Inches(5.8), Inches(0.38),
    size=13, bold=True, color=RED)

survival_table = [
    ("Off-by-one arithmetic  (+1/−1)", "38",  "Assertions too loose"),
    ("Comparison flips  (>= vs >)",    "31",  "Not on critical path"),
    ("String literal changes",         "28",  "Not asserted"),
    ("Dead / unreachable code",        "22",  "Cannot be killed"),
    ("Boolean negation",               "15",  "Equivalent mutant"),
    ("Other",                          "10",  "Various"),
]
box(sl, Inches(0.4), Inches(3.76), Inches(6.1), Inches(0.3),
    RGBColor(0x1A, 0x22, 0x2E))
txt(sl, "Mutant Type", Inches(0.5), Inches(3.78), Inches(3.5), Inches(0.25), size=9, bold=True, color=MUTED)
txt(sl, "#",           Inches(4.1), Inches(3.78), Inches(0.5), Inches(0.25), size=9, bold=True, color=MUTED)
txt(sl, "Reason",      Inches(4.7), Inches(3.78), Inches(1.6), Inches(0.25), size=9, bold=True, color=MUTED)
for i, (typ, cnt, reason) in enumerate(survival_table):
    y = Inches(4.08) + i * Inches(0.5)
    bg = CARD2 if i % 2 == 0 else SURFACE
    box(sl, Inches(0.4), y, Inches(6.1), Inches(0.46), bg)
    txt(sl, typ,    Inches(0.5), y + Inches(0.07), Inches(3.5), Inches(0.32), size=10, color=WHITE)
    txt(sl, cnt,    Inches(4.1), y + Inches(0.07), Inches(0.5), Inches(0.32), size=11, bold=True, color=RED)
    txt(sl, reason, Inches(4.7), y + Inches(0.07), Inches(1.6), Inches(0.32), size=9,  color=MUTED)

# Right — mutation-targeted tests
card(sl, Inches(6.75), Inches(3.2), Inches(6.25), Inches(3.93))
txt(sl, "Our Response — 73 Precision Tests", Inches(6.95), Inches(3.32), Inches(5.9), Inches(0.38),
    size=13, bold=True, color=ORANGE)
txt(sl, "Each test targets exactly one surviving mutant.\nDocstrings use labels like [DAT-L234] / [CAL-L78].",
    Inches(6.95), Inches(3.76), Inches(5.9), Inches(0.55),
    size=11, color=MUTED)

mutation_classes = [
    ("[DAT-L234]", "Timer parity edge",      "6 tests"),
    ("[DAT-L89]",  "Collection guard",       "8 tests"),
    ("[DAT-L312]", "Subtask prefix",         "5 tests"),
    ("[CAL-L78]",  "Leap year logic",        "9 tests"),
    ("[CAL-L145]", "CSV row detection",      "7 tests"),
    ("[CAL-L201]", "ICS RRULE parsing",      "6 tests"),
    ("[CAL-L310]", "Year filter boundary",   "8 tests"),
    ("Other",      "Various data mutations", "24 tests"),
]
for i, (tag, desc, count) in enumerate(mutation_classes):
    y = Inches(4.42) + i * Inches(0.45)
    bg = CARD2 if i % 2 == 0 else SURFACE
    box(sl, Inches(6.8), y, Inches(6.15), Inches(0.41), bg)
    txt(sl, tag,   Inches(6.9),  y + Inches(0.06), Inches(1.3), Inches(0.3), size=10, bold=True, color=TEAL)
    txt(sl, desc,  Inches(8.3),  y + Inches(0.06), Inches(3.0), Inches(0.3), size=10, color=WHITE)
    txt(sl, count, Inches(11.4), y + Inches(0.06), Inches(1.3), Inches(0.3), size=10, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CLI + MOCK + INTEGRATION
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "CLI · Mock · Integration Testing", "3 additional suites — 65 tests total", TEAL)

suites_data = [
    ("CLI Testing", "27 Tests", ACCENT,
     "Two-tier approach needed because Calcure requires a real TTY (curses).",
     [("Tier 1 — no TTY", "subprocess", "-v, --folder, --config flags"),
      ("Tier 2 — PTY",    "pty module", "--task, --event (curses UI)"),
      ("Exit codes",      "verified",   "0 = success, non-0 = error"),
      ("Output patterns", "regex match","version string, help text")],
     "sys.executable used — avoids hardcoded Python paths"),
    ("Mock Testing", "17 Tests", ORANGE,
     "Mocking replaces slow/side-effectful dependencies during unit tests.",
     [("time.time()",         "unittest.mock", "Timer parity tests"),
      ("builtins.open",       "mock_open",     "File I/O isolation"),
      ("urllib.request",      "MagicMock",     "ICS download tests"),
      ("pathlib.Path.replace","patch",         "Atomic save tests"),
      ("datetime.now()",      "patch",         "Date-sensitive tests")],
     "Verified call counts + return values in every test"),
    ("Integration Testing", "21 Tests", PURPLE,
     "Tests data flow through multiple connected components end-to-end.",
     [("CSV round-trip",     "write → read",  "Data survives file I/O"),
      ("ICS import chain",   "parse → store", "Events stored correctly"),
      ("Task pipeline",      "add → save → load","Full lifecycle"),
      ("Config + loader",    "config → parse","Settings applied"),
      ("Multi-event search", "filter → render","Combined operation")],
     "Catches bugs that unit tests miss — especially format issues"),
]
for i, (title, count, color, desc, table, note) in enumerate(suites_data):
    x = Inches(0.35) + i * Inches(4.35)
    card(sl, x, Inches(1.35), Inches(4.1), Inches(5.78))
    box(sl, x, Inches(1.35), Inches(4.1), Inches(0.07), color)
    pill(sl, count, x + Inches(0.18), Inches(1.52), Inches(1.4), Inches(0.34), color, size=11)
    txt(sl, title, x + Inches(0.18), Inches(1.96), Inches(3.7), Inches(0.42),
        size=15, bold=True, color=WHITE)
    txt(sl, desc, x + Inches(0.18), Inches(2.44), Inches(3.7), Inches(0.55),
        size=10, color=MUTED, italic=True)
    # mini table
    box(sl, x + Inches(0.12), Inches(3.08), Inches(3.85), Inches(0.28),
        RGBColor(0x1A, 0x22, 0x2E))
    txt(sl, "What",   x + Inches(0.2),  Inches(3.1), Inches(1.4), Inches(0.24), size=8, bold=True, color=MUTED)
    txt(sl, "How",    x + Inches(1.62), Inches(3.1), Inches(1.1), Inches(0.24), size=8, bold=True, color=MUTED)
    txt(sl, "Checks", x + Inches(2.75), Inches(3.1), Inches(1.1), Inches(0.24), size=8, bold=True, color=MUTED)
    for j, (what, how, checks) in enumerate(table):
        y = Inches(3.38) + j * Inches(0.42)
        bg = CARD2 if j % 2 == 0 else SURFACE
        box(sl, x + Inches(0.12), y, Inches(3.85), Inches(0.38), bg)
        txt(sl, what,   x + Inches(0.2),  y + Inches(0.05), Inches(1.4), Inches(0.3), size=9, color=color)
        txt(sl, how,    x + Inches(1.62), y + Inches(0.05), Inches(1.1), Inches(0.3), size=9, color=WHITE)
        txt(sl, checks, x + Inches(2.75), y + Inches(0.05), Inches(1.1), Inches(0.3), size=8, color=MUTED)
    # note
    txt(sl, f"★ {note}", x + Inches(0.18), Inches(5.6), Inches(3.7), Inches(0.42),
        size=9, color=GREEN, italic=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BUG 1
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Bug #1 — Whitespace-Only Name Accepted", "Discovered via Error Guessing · data.py", RED)

# Bug info card
card(sl, Inches(0.35), Inches(1.35), Inches(12.63), Inches(5.78))
box(sl, Inches(0.35), Inches(1.35), Inches(0.10), Inches(5.78), RED)

txt(sl, "Description", Inches(0.6), Inches(1.5), Inches(12.0), Inches(0.38),
    size=12, bold=True, color=MUTED)
txt(sl, 'A task with name = "   " (three spaces) should be rejected like an empty name. '
        'Instead, it is silently stored — the guard only checks len(name) > 0, '
        'and len("   ") == 3 so it passes.',
    Inches(0.6), Inches(1.92), Inches(12.0), Inches(0.72),
    size=13, color=WHITE)

# Two code boxes side by side
txt(sl, "Buggy Code  (data.py ~L89)", Inches(0.6), Inches(2.78), Inches(5.8), Inches(0.35),
    size=11, bold=True, color=RED)
code_box(sl,
    'def add_item(self, item):\n'
    '    if len(item.name) > 0:   # BUG\n'
    '        self.items.append(item)\n'
    '    # "   " has len 3 — passes!',
    Inches(0.6), Inches(3.15), Inches(5.8), Inches(1.08), RED)

txt(sl, "Suggested Fix", Inches(7.0), Inches(2.78), Inches(5.8), Inches(0.35),
    size=11, bold=True, color=GREEN)
code_box(sl,
    'def add_item(self, item):\n'
    '    if len(item.name.strip()) > 0:  # FIX\n'
    '        self.items.append(item)\n'
    '    # "   ".strip() == "" — rejected ✓',
    Inches(7.0), Inches(3.15), Inches(5.8), Inches(1.08), GREEN)

# Impact table
txt(sl, "Impact Analysis", Inches(0.6), Inches(4.38), Inches(12.0), Inches(0.35),
    size=12, bold=True, color=ORANGE)
impacts = [
    ("Symptom",    "Tasks with blank names stored and appear as empty rows in the UI"),
    ("Severity",   "Medium — UI confusion, data corruption in exported CSV"),
    ("Detected by","test_add_task_whitespace_only_name_rejected  [EG blackbox test]"),
    ("Fix effort", "One-line — add .strip() before len() check"),
]
for i, (key, val) in enumerate(impacts):
    y = Inches(4.82) + i * Inches(0.52)
    bg = CARD2 if i % 2 == 0 else SURFACE
    box(sl, Inches(0.5), y, Inches(12.4), Inches(0.48), bg)
    txt(sl, key,  Inches(0.65), y + Inches(0.1), Inches(1.6), Inches(0.3), size=11, bold=True, color=ORANGE)
    txt(sl, val,  Inches(2.35), y + Inches(0.1), Inches(10.4), Inches(0.3), size=11, color=WHITE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — BUG 2
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Bug #2 — CSV Crash on Malformed Year Field", "Discovered via Error Guessing · loaders.py", RED)

card(sl, Inches(0.35), Inches(1.35), Inches(12.63), Inches(5.78))
box(sl, Inches(0.35), Inches(1.35), Inches(0.10), Inches(5.78), RED)

txt(sl, "Description", Inches(0.6), Inches(1.5), Inches(12.0), Inches(0.38),
    size=12, bold=True, color=MUTED)
txt(sl, 'loaders.py calls int(row[1]) to parse the year column without any error handling. '
        'If a CSV file contains a corrupted row like  0,not-a-year,9,15,"Conference",3,weekly,normal  '
        'the application raises ValueError and crashes at startup — showing no error message.',
    Inches(0.6), Inches(1.92), Inches(12.0), Inches(0.72),
    size=13, color=WHITE)

txt(sl, "Buggy Code  (loaders.py ~L145)", Inches(0.6), Inches(2.78), Inches(5.8), Inches(0.35),
    size=11, bold=True, color=RED)
code_box(sl,
    'for row in reader:\n'
    '    year  = int(row[1])   # ValueError if bad\n'
    '    month = int(row[2])\n'
    '    day   = int(row[3])',
    Inches(0.6), Inches(3.15), Inches(5.8), Inches(1.08), RED)

txt(sl, "Suggested Fix", Inches(7.0), Inches(2.78), Inches(5.8), Inches(0.35),
    size=11, bold=True, color=GREEN)
code_box(sl,
    'for row in reader:\n'
    '    try:\n'
    '        year  = int(row[1])\n'
    '        month = int(row[2])\n'
    '        day   = int(row[3])\n'
    '    except (ValueError, IndexError):\n'
    '        continue   # skip bad rows',
    Inches(7.0), Inches(3.15), Inches(5.8), Inches(1.08), GREEN)

impacts = [
    ("Symptom",    "Application crashes silently on startup if any CSV row has a non-integer year"),
    ("Severity",   "High — data corruption or accidental edits render the app unlaunchable"),
    ("Detected by","test_event_loader_malformed_csv_does_not_crash  [EG blackbox test]"),
    ("Fix effort", "Wrap int() casts in try/except; continue to skip malformed rows"),
]
txt(sl, "Impact Analysis", Inches(0.6), Inches(4.38), Inches(12.0), Inches(0.35),
    size=12, bold=True, color=ORANGE)
for i, (key, val) in enumerate(impacts):
    y = Inches(4.82) + i * Inches(0.52)
    bg = CARD2 if i % 2 == 0 else SURFACE
    box(sl, Inches(0.5), y, Inches(12.4), Inches(0.48), bg)
    txt(sl, key,  Inches(0.65), y + Inches(0.1), Inches(1.6), Inches(0.3), size=11, bold=True, color=ORANGE)
    txt(sl, val,  Inches(2.35), y + Inches(0.1), Inches(10.4), Inches(0.3), size=11, color=WHITE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTS DASHBOARD
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Results Dashboard", "Final numbers at a glance", GREEN)

# 6 big metrics top row
bigs = [
    ("341",   "Total Tests",       GREEN,  "across 6 suites"),
    ("89%",   "Branch Coverage",   ACCENT, "core modules"),
    ("81.6%", "Mutation Score",    ORANGE, "above 80% target"),
    ("2",     "Real Bugs Found",   RED,    "with root causes"),
    ("6",     "Test Suites",       PURPLE, "complementary"),
    ("5,784", "LOC Under Test",    TEAL,   "Python 3.9+"),
]
cols3 = 3
for i, (val, lbl, col, sub) in enumerate(bigs):
    c = i % cols3
    r = i // cols3
    x = Inches(0.35) + c * Inches(4.35)
    y = Inches(1.35) + r * Inches(2.15)
    metric_card(sl, x, y, Inches(4.0), Inches(2.0), val, lbl, col, sub)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BEFORE vs AFTER
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Before vs After", "What we built over the semester", ACCENT)

befores = [
    ("Test Files",       "2",   "6",     "+4 new suites"),
    ("Total Tests",      "5",   "341",   "+336 tests"),
    ("Branch Coverage",  "~0%", "89%",   "+89 percentage points"),
    ("Mutation Score",   "0%",  "81.6%", "above 80% industry target"),
    ("Bugs Found",       "0",   "2",     "both with root-cause analysis"),
    ("Lines of Tests",   "67",  "4,500+","×67 growth"),
    ("Techniques Used",  "0",   "7",     "EP·BA·EG·CT·WB·Mock·Mutation"),
    ("Tools Used",       "0",   "5",     "pytest·coverage·mutmut·pty·mock"),
]

# Header row
col_xs_ba = [Inches(0.35), Inches(4.3), Inches(6.5), Inches(8.9)]
col_ws_ba = [Inches(3.85), Inches(2.1), Inches(2.3), Inches(4.0)]
box(sl, Inches(0.35), Inches(1.35), SW - Inches(0.7), Inches(0.42),
    RGBColor(0x1A, 0x22, 0x2E))
for j, h in enumerate(["Metric", "Before", "After", "Delta"]):
    align = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    col = RED if j == 1 else (GREEN if j == 2 else (ACCENT if j == 3 else MUTED))
    txt(sl, h, col_xs_ba[j], Inches(1.38), col_ws_ba[j], Inches(0.35),
        size=11, bold=True, color=col, align=align)

for i, (metric, before, after, delta) in enumerate(befores):
    y = Inches(1.77) + i * Inches(0.68)
    bg = SURFACE if i % 2 == 0 else BG
    box(sl, Inches(0.35), y, SW - Inches(0.7), Inches(0.64), bg)
    txt(sl, metric, col_xs_ba[0], y + Inches(0.12), col_ws_ba[0], Inches(0.42), size=13, color=WHITE)
    txt(sl, before, col_xs_ba[1], y + Inches(0.08), col_ws_ba[1], Inches(0.5),
        size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)
    txt(sl, "→", Inches(6.3), y + Inches(0.12), Inches(0.5), Inches(0.38),
        size=16, color=DIM, align=PP_ALIGN.CENTER)
    txt(sl, after,  col_xs_ba[2], y + Inches(0.08), col_ws_ba[2], Inches(0.5),
        size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(sl, delta,  col_xs_ba[3], y + Inches(0.12), col_ws_ba[3], Inches(0.42),
        size=11, color=MUTED, italic=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — LESSONS LEARNED
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Lessons Learned", "Key takeaways from testing Calcure")

lessons = [
    (RED,    "Error Guessing reveals real bugs that EP & BA miss",
     "Both bugs were found by asking 'what if the data is bad?' — not by partitioning valid input. Always include EG."),
    (GREEN,  "High branch coverage does not mean high confidence",
     "89% branch coverage still left 144 surviving mutants. Coverage proves execution, not verification."),
    (ORANGE, "Mutation score exposes weak assertions precisely",
     "Each surviving mutant pointed to a test whose assertion was too broad. Precision tests are more valuable than many weak ones."),
    (PURPLE, "Mocks are fast but integration tests are essential",
     "Mocks isolated components well, but integration tests caught CSV round-trip bugs that mocks concealed."),
    (TEAL,   "TUI apps require a PTY — standard subprocess is not enough",
     "Calcure needs a real TTY. Without the pty-based Tier 2 tests, 40% of CLI flags would be completely untestable."),
]
for i, (color, title, body) in enumerate(lessons):
    y = Inches(1.42) + i * Inches(1.14)
    card(sl, Inches(0.35), y, SW - Inches(0.7), Inches(1.05))
    box(sl, Inches(0.35), y, Inches(0.10), Inches(1.05), color)
    # number badge
    box(sl, Inches(0.6), y + Inches(0.28), Inches(0.42), Inches(0.42), color)
    txt(sl, str(i+1), Inches(0.6), y + Inches(0.28), Inches(0.42), Inches(0.42),
        size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(sl, title, Inches(1.2), y + Inches(0.1), Inches(11.5), Inches(0.42),
        size=13, bold=True, color=WHITE)
    txt(sl, body, Inches(1.2), y + Inches(0.55), Inches(11.5), Inches(0.42),
        size=11, color=MUTED)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DEMO / REPO
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
header(sl, "Demo & Repository", "Run the tests yourself", TEAL)

card(sl, Inches(0.35), Inches(1.35), Inches(12.63), Inches(1.55))
txt(sl, "GitHub Repository", Inches(0.55), Inches(1.48), Inches(5.0), Inches(0.38),
    size=12, bold=True, color=TEAL)
txt(sl, "https://github.com/Yashcoder2802/calcure-testing",
    Inches(0.55), Inches(1.9), Inches(10.0), Inches(0.42),
    size=16, bold=True, color=WHITE)
txt(sl, "Branch: main   |   Tests: tests/   |   Report: TESTING_REPORT.md   |   Slides: reports/",
    Inches(0.55), Inches(2.36), Inches(12.0), Inches(0.38),
    size=11, color=MUTED)

commands = [
    ("Run all tests",            "python -m pytest tests/ -v",                    ACCENT),
    ("Branch coverage report",   "coverage run -m pytest tests/ && coverage report --branch", GREEN),
    ("Mutation testing",         "mutmut run && mutmut results",                  ORANGE),
    ("Specific suite",           "python -m pytest tests/blackbox/ -v",          PURPLE),
    ("HTML coverage report",     "coverage html && open htmlcov/index.html",     TEAL),
]
card(sl, Inches(0.35), Inches(3.1), Inches(12.63), Inches(3.92))
txt(sl, "Quick Start Commands", Inches(0.55), Inches(3.22), Inches(6.0), Inches(0.38),
    size=13, bold=True, color=TEAL)
box(sl, Inches(0.35), Inches(3.65), Inches(12.63), Inches(0.3),
    RGBColor(0x1A, 0x22, 0x2E))
txt(sl, "Description", Inches(0.55), Inches(3.67), Inches(3.5), Inches(0.25), size=9, bold=True, color=MUTED)
txt(sl, "Command",     Inches(4.2),  Inches(3.67), Inches(8.6), Inches(0.25), size=9, bold=True, color=MUTED)
for i, (desc, cmd, color) in enumerate(commands):
    y = Inches(3.97) + i * Inches(0.56)
    bg = CARD2 if i % 2 == 0 else SURFACE
    box(sl, Inches(0.4), y, Inches(12.55), Inches(0.52), bg)
    txt(sl, desc, Inches(0.55), y + Inches(0.1), Inches(3.5), Inches(0.35), size=11, color=WHITE)
    txt(sl, cmd,  Inches(4.2),  y + Inches(0.1), Inches(8.6), Inches(0.35), size=11, bold=True, color=color)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Q&A
# ═════════════════════════════════════════════════════════════════════════════
sl = add_slide()
# Split background
box(sl, 0, 0, SW, SH, RGBColor(0x0A, 0x0E, 0x14))
box(sl, 0, 0, SW, Inches(0.08), ACCENT)
box(sl, 0, SH - Inches(0.08), SW, Inches(0.08), ACCENT)

# Giant Q&A text
txt(sl, "Q&A", Inches(0), Inches(1.4), SW, Inches(2.8),
    size=120, bold=True, color=RGBColor(0x1A, 0x22, 0x2E), align=PP_ALIGN.CENTER)
# Overlay slightly lighter
txt(sl, "Q&A", Inches(0.05), Inches(1.35), SW, Inches(2.8),
    size=120, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Thank you for listening!", Inches(0), Inches(4.3), SW, Inches(0.65),
    size=24, color=ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "Yash Salunke  ·  Harshal Gajjar",
    Inches(0), Inches(4.98), SW, Inches(0.45),
    size=15, color=MUTED, align=PP_ALIGN.CENTER)

# summary pills row
for i, (label, color) in enumerate([
    ("341 Tests", GREEN), ("89% Coverage", ACCENT),
    ("81.6% Mutation", ORANGE), ("2 Bugs Found", RED),
    ("6 Suites", PURPLE),
]):
    pill(sl, label, Inches(1.0 + i * 2.28), Inches(5.88), Inches(2.1), Inches(0.46), color, size=12)

# ─────────────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "presentation.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
